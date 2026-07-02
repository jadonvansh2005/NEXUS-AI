from typing import List

from pptx import Presentation

from rag.loaders.base_loader import (
    BaseLoader
)

from rag.loaders.loader_constants import (
    PPTX_EXTENSIONS,
    DocumentType
)

from rag.loaders.loader_factory import (
    LoaderFactory
)

from rag.loaders.loader_schema import (
    LoadedDocument
)

from rag.loaders.loader_utils import (
    LoaderUtils
)


class PPTLoader(

    BaseLoader

):

    SUPPORTED_EXTENSIONS = (

        PPTX_EXTENSIONS

    )

    def load(

        self

    ) -> LoadedDocument:

        presentation = (

            Presentation(

                self.absolute_path

            )

        )

        slides: List[str] = []

        for slide in (

            presentation.slides

        ):

            slide_text = []

            for shape in (

                slide.shapes

            ):

                if hasattr(

                    shape,

                    "text"

                ):

                    text = (

                        shape.text

                    ).strip()

                    if text:

                        slide_text.append(

                            text

                        )

            slides.append(

                LoaderUtils.normalize_text(

                    "\n".join(

                        slide_text

                    )

                )

            )

        full_text = (

            "\n\n".join(

                slides

            )

        )

        return LoadedDocument(

            text=full_text,

            document_type=(

                DocumentType.PPTX.value

            ),

            source_path=(

                self.absolute_path

            ),

            filename=(

                self.filename

            ),

            metadata=self.metadata(),

            page_count=len(

                slides

            ),

            pages=slides,

            checksum=(

                LoaderUtils.calculate_checksum(

                    self.absolute_path

                )

            ),

            file_size=(

                self.filesize

            )

        )


LoaderFactory.register(

    PPTLoader

)